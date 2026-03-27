"""PowerPoint mixin class."""

from __future__ import annotations

import io
import typing as ty
from contextlib import contextmanager
from enum import IntEnum
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from loguru import logger

from koyo.typing import PathLike
from koyo.utilities import is_installed

HAS_PPTX = is_installed("pptx")

if ty.TYPE_CHECKING:
    from PIL import Image

    try:
        from pptx.presentation import Presentation
    except ImportError:
        Presentation = None  # type: ignore[assignment,misc]


DEFAULT_PIL_FMT = "PNG"
DEFAULT_MPL_FMT = "PNG"


class SlideLayout(IntEnum):
    """Slide layout options for PowerPoint presentations.

    Each member corresponds to a slide layout index in the default
    python-pptx template.
    """

    TITLE = 0
    TITLE_AND_CONTENT = 1
    SECTION_HEADER = 2
    TITLE_AND_TWO_CONTENT = 3
    TITLE_AND_TWO_CONTENT_WITH_HEADER = 4
    TITLE_AND_BLANK = 5
    BLANK = 6
    HEADER_AND_TWO_CONTENT = 7
    PICTURE_AND_CAPTION = 8


class PPTXMixin:
    """Mixin class to help export figures to PPTX files.

    Classes that inherit from this mixin gain methods for creating and
    populating a PowerPoint presentation with matplotlib figures or PIL
    images. Subclasses must implement `pptx_filename`.

    Attributes
    ----------
    _pptx : Presentation or None
        Cached presentation object. Initialized lazily when first accessed.
    as_pptx : bool
        Whether to export output as a PPTX file instead of saving individual
        image files.
    """

    _pptx = None
    as_pptx: bool = False

    @property
    def pptx_filename(self) -> Path:
        """Return the path to the PPTX output file.

        Returns
        -------
        Path
            Destination path for the PPTX file.

        Raises
        ------
        NotImplementedError
            Always — subclasses must override this property.
        """
        raise NotImplementedError("Must implement method")

    @property
    def pptx(self) -> Presentation | None:
        """Return the current PowerPoint presentation object.

        Lazily creates a new presentation at `pptx_filename` if `as_pptx`
        is ``True`` and no presentation has been initialized yet.

        Returns
        -------
        Presentation or None
            The active presentation, or ``None`` if `as_pptx` is ``False``.
        """
        if self._pptx is None and self.as_pptx:
            self._pptx = self._make_pptx(self.pptx_filename)
        return self._pptx

    @pptx.setter
    def pptx(self, value: Presentation | None) -> None:
        """Set the current PowerPoint presentation, saving any existing one first.

        If a presentation is already cached, it is saved to `pptx_filename`
        before being replaced.

        Parameters
        ----------
        value : Presentation or None
            The new presentation to cache, or ``None`` to clear the cache.
        """
        if self._pptx is not None:
            self._pptx.save(self.pptx_filename)  # type: ignore[arg-type]
        self._pptx = value

    @staticmethod
    def _make_pptx(filename: PathLike) -> Presentation:
        """Create a new widescreen (16x9) PowerPoint presentation.

        Parameters
        ----------
        filename : str or Path
            Path at which the presentation will eventually be saved. Stored
            as a ``_filename`` attribute on the returned object.

        Returns
        -------
        Presentation
            A new, empty python-pptx ``Presentation`` sized 16 x 9 inches.

        Raises
        ------
        ImportError
            If ``python-pptx`` is not installed.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            raise ImportError(
                "pptx is not installed. Please install it using `pip install python-pptx`.",
            ) from None

        pptx = Presentation()
        pptx.slide_width = Inches(16)
        pptx.slide_height = Inches(9)
        pptx._filename = filename  # type: ignore[attr-defined]
        return pptx

    @contextmanager
    def _export_pptx_figures(
        self,
        filename: PathLike | None = None,
    ) -> ty.Generator[Presentation | None, None, None]:
        """Context manager that yields a presentation for bulk figure export.

        Switches matplotlib to the non-interactive ``agg`` backend, yields
        the active (or newly created) presentation, and saves it on exit.
        Any exception that occurs inside the block is logged rather than
        re-raised, so that partial results are still saved.

        Parameters
        ----------
        filename : str or Path, optional
            If provided, a fresh presentation is created at this path instead
            of reusing the cached one. When ``None``, the cached presentation
            (from ``self.pptx``) is used and reset to ``None`` after saving.

        Yields
        ------
        Presentation or None
            The presentation to populate, or ``None`` if `as_pptx` is
            ``False``.
        """
        import matplotlib as mpl

        mpl.use("agg")

        pptx: Presentation | None = None
        reset = False
        if self.as_pptx:
            if filename:
                pptx = self._make_pptx(filename)
            else:
                pptx = self.pptx
                reset = True
        try:
            yield pptx
        except Exception as exc:
            logger.exception(f"Error exporting to PPTX: {exc}")
        finally:
            if self.as_pptx:
                self._save_pptx(pptx, filename, reset)  # type: ignore[arg-type]

    def _add_title_to_pptx(self, title: str, pptx: Presentation | None = None) -> None:
        """Add a title slide to the presentation.

        Parameters
        ----------
        title : str
            Text to display on the title slide. Newlines are supported and
            will cause the title box to expand proportionally.
        pptx : Presentation, optional
            Target presentation. Defaults to ``self.pptx``.
        """
        pptx = pptx or self.pptx
        add_title_to_pptx(pptx, title)  # type: ignore[arg-type]

    def _add_content_to_pptx(
        self,
        content: str,
        title: str = "",
        pptx: Presentation | None = None,
    ) -> None:
        """Add a title-and-content slide to the presentation.

        Parameters
        ----------
        content : str
            Body text for the content placeholder.
        title : str, optional
            Slide title. Defaults to an empty string (no title text).
        pptx : Presentation, optional
            Target presentation. Defaults to ``self.pptx``.
        """
        pptx = pptx or self.pptx
        add_content_to_pptx(pptx, content, title)  # type: ignore[arg-type]

    def _add_mpl_figure_to_pptx(
        self,
        filename: PathLike,
        fig: plt.Figure,
        face_color: str | np.ndarray | None = None,
        bbox_inches: str | None = "tight",
        dpi: int = 150,
        override: bool = False,
        if_empty: str = "warn",
        close: bool = False,
        title: str = "",
        pptx: Presentation | None = None,
        **kwargs: ty.Any,
    ) -> None:
        """Add a matplotlib figure as a slide, or save it to disk.

        When a presentation is available the figure is rendered into an
        in-memory buffer and inserted as a picture on a new slide. Otherwise
        it is saved directly to `filename`.

        Parameters
        ----------
        filename : str or Path
            Fallback output path used when `pptx` is ``None``.
        fig : matplotlib.figure.Figure
            The figure to export. If ``None``, the `if_empty` handler is
            invoked and the method returns early.
        face_color : str or numpy.ndarray, optional
            Background color passed to ``fig.savefig``. Defaults to the
            figure's own face color when ``None``.
        bbox_inches : str, optional
            Bounding-box argument forwarded to ``fig.savefig``. Defaults to
            ``"tight"``.
        dpi : int, optional
            Resolution in dots per inch. Defaults to ``150``.
        override : bool, optional
            When ``True``, overwrite `filename` even if it already exists.
            Ignored when writing to a presentation. Defaults to ``False``.
        if_empty : str, optional
            Behaviour when `fig` is ``None``: ``"none"`` silently returns,
            ``"warn"`` logs a warning, ``"raise"`` raises a ``ValueError``.
            Defaults to ``"warn"``.
        close : bool, optional
            Whether to close the figure after export. Defaults to ``False``.
        title : str, optional
            Slide title. Defaults to an empty string.
        pptx : Presentation, optional
            Target presentation. Defaults to ``self.pptx``.
        **kwargs
            Additional keyword arguments forwarded to ``fig.savefig``.
        """
        if fig is None:
            self._inform_on_empty(if_empty)
            return

        pptx = pptx or self.pptx
        add_mpl_figure_to_pptx(
            pptx,
            filename,
            fig,
            face_color,
            bbox_inches,
            dpi,
            override,
            close=close,
            title=title,
            **kwargs,
        )

    def _add_pil_image_to_pptx(
        self,
        filename: PathLike,
        image: Image,
        dpi: int = 150,
        fmt: str = "JPEG",
        override: bool = False,
        close: bool = False,
        title: str = "",
        pptx: Presentation | None = None,
        **kwargs: ty.Any,
    ) -> None:
        """Add a PIL image as a slide, or save it to disk.

        When a presentation is available the image is serialized into an
        in-memory buffer and inserted as a picture on a new slide. Otherwise
        it is saved directly to `filename`.

        Parameters
        ----------
        filename : str or Path
            Fallback output path used when `pptx` is ``None``.
        image : PIL.Image.Image
            The image to export.
        dpi : int, optional
            Resolution in dots per inch. Defaults to ``150``.
        fmt : str, optional
            PIL format string (e.g. ``"JPEG"``, ``"PNG"``). Defaults to
            ``"JPEG"``.
        override : bool, optional
            When ``True``, overwrite `filename` even if it already exists.
            Ignored when writing to a presentation. Defaults to ``False``.
        close : bool, optional
            Whether to close the image after export. Defaults to ``False``.
        title : str, optional
            Slide title. Defaults to an empty string.
        pptx : Presentation, optional
            Target presentation. Defaults to ``self.pptx``.
        **kwargs
            Additional keyword arguments forwarded to ``image.save``.
        """
        pptx = pptx or self.pptx
        add_pil_image_to_pptx(
            pptx,
            filename,
            image,
            dpi,
            fmt=fmt,
            override=override,
            close=close,
            title=title,
            **kwargs,
        )

    def _save_pptx(
        self,
        pptx: Presentation,
        filename: PathLike | None = None,
        reset: bool = False,
    ) -> None:
        """Save the presentation to disk.

        The save path is resolved in the following priority order:

        1. ``pptx._filename`` if that attribute exists (set by ``_make_pptx``).
        2. The `filename` argument.
        3. ``self.pptx_filename``.

        Parameters
        ----------
        pptx : Presentation
            The presentation to save.
        filename : str or Path, optional
            Explicit destination path. Used only when ``pptx._filename`` is
            not set.
        reset : bool, optional
            When ``True``, clears the cached ``self._pptx`` after saving.
            Defaults to ``False``.
        """
        filename = pptx._filename if hasattr(pptx, "_filename") else filename or self.pptx_filename
        if len(pptx.slides) == 0:
            logger.warning("No slides to save in PPTX")
        else:
            pptx.save(filename)  # type: ignore[arg-type]
            logger.trace(f"Saved PPTX to {filename} with {len(pptx.slides)} slides")
        if reset:
            self._pptx = None

    @staticmethod
    def _inform_on_empty(if_empty: str = "warn") -> None:
        """Handle the case where a figure or image is ``None``.

        Parameters
        ----------
        if_empty : str, optional
            Controls the behaviour:

            - ``"none"`` — do nothing.
            - ``"warn"`` — emit a loguru warning (default).
            - ``"raise"`` — raise a ``ValueError``.

        Raises
        ------
        ValueError
            If `if_empty` is ``"raise"``.
        """
        if if_empty == "none":
            return
        if if_empty == "warn":
            logger.warning("Figure was empty")
        elif if_empty == "raise":
            raise ValueError("Figure was empty")


def add_title_to_pptx(pptx: Presentation, title: str) -> None:
    """Add a title slide to a PowerPoint presentation.

    The title text box is resized vertically to accommodate multi-line titles,
    and stretched to the full slide width.

    Parameters
    ----------
    pptx : Presentation
        The presentation to which the slide is appended.
    title : str
        Title text. Newline characters cause the title box to grow
        proportionally.
    """
    from pptx.util import Cm

    slide = pptx.slides.add_slide(pptx.slide_layouts[SlideLayout.TITLE])
    line_count = title.count("\n") + 1
    height = slide.shapes.title.height * line_count
    slide.shapes.title.text = title
    slide.shapes.title.left = Cm(0)
    slide.shapes.title.width = pptx.slide_width
    slide.shapes.title.height = height


def add_content_to_pptx(pptx: Presentation, content: str, title: str = "") -> None:
    """Add a title-and-content slide to a PowerPoint presentation.

    The title text box is resized vertically to accommodate multi-line titles
    and stretched to the full slide width. The content placeholder receives
    `content` verbatim.

    Parameters
    ----------
    pptx : Presentation
        The presentation to which the slide is appended.
    content : str
        Body text placed in the content placeholder.
    title : str, optional
        Slide title. Defaults to an empty string.
    """
    from pptx.util import Cm

    slide = pptx.slides.add_slide(pptx.slide_layouts[SlideLayout.TITLE_AND_CONTENT])
    line_count = title.count("\n") + 1
    height = slide.shapes.title.height * line_count
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    slide.shapes.title.left = Cm(0)
    slide.shapes.title.width = pptx.slide_width
    slide.shapes.title.height = height


def add_mpl_figure_to_pptx(
    pptx: Presentation | None,
    filename: PathLike,
    fig: plt.Figure,
    face_color: str | np.ndarray | None = None,
    bbox_inches: str | None = "tight",
    dpi: int = 150,
    override: bool = False,
    close: bool = False,
    title: str = "",
    image_format: str = DEFAULT_MPL_FMT,
    **kwargs: ty.Any,
) -> None:
    """Add a matplotlib figure to a presentation slide, or save it to disk.

    When `pptx` is provided the figure is rendered into an in-memory PNG (or
    other `image_format`) buffer and inserted as a picture on a new slide.
    When `pptx` is ``None`` the figure is saved to `filename` instead
    (subject to the `override` flag).

    Parameters
    ----------
    pptx : Presentation or None
        Target presentation. When ``None``, the figure is saved to disk.
    filename : str or Path
        Destination path used when `pptx` is ``None``.
    fig : matplotlib.figure.Figure
        The figure to export.
    face_color : str or numpy.ndarray, optional
        Background color for the exported image. Defaults to the figure's
        own face color when ``None``.
    bbox_inches : str, optional
        Bounding-box argument forwarded to ``fig.savefig``. Defaults to
        ``"tight"``.
    dpi : int, optional
        Resolution in dots per inch. Defaults to ``150``.
    override : bool, optional
        When ``True``, overwrite `filename` even if it already exists.
        Ignored when `pptx` is not ``None``. Defaults to ``False``.
    close : bool, optional
        Whether to close the figure after export. Defaults to ``False``.
    title : str, optional
        Slide title inserted above the picture. Defaults to an empty string.
    image_format : str, optional
        Image format string passed to ``fig.savefig`` (e.g. ``"PNG"``).
        Defaults to :data:`DEFAULT_MPL_FMT`.
    **kwargs
        Additional keyword arguments forwarded to ``fig.savefig``.
    """
    face_color = face_color if face_color is not None else fig.get_facecolor()
    if pptx is not None:
        with io.BytesIO() as image_stream:
            fig.savefig(
                image_stream,
                dpi=dpi,
                facecolor=face_color,
                bbox_inches=bbox_inches,
                format=image_format,
                **kwargs,
            )
            slide, left, top = _insert_slide(pptx, title=title)
            slide.shapes.add_picture(image_stream, left, top)
    else:
        if override or not Path(filename).exists():
            fig.savefig(filename, dpi=dpi, facecolor=face_color, bbox_inches=bbox_inches, **kwargs)
    if close:
        plt.close(fig)


def add_pil_image_to_pptx(
    pptx: Presentation | None,
    filename: PathLike,
    image: Image,
    dpi: int = 150,
    fmt: str = DEFAULT_PIL_FMT,
    override: bool = False,
    close: bool = False,
    title: str = "",
    **kwargs: ty.Any,
) -> None:
    """Add a PIL image to a presentation slide, or save it to disk.

    When `pptx` is provided the image is serialized into an in-memory buffer
    and inserted as a picture on a new slide. RGBA images are converted to
    RGB automatically when `fmt` is ``"JPEG"`` or ``"JPG"`` (which do not
    support an alpha channel). When `pptx` is ``None`` the image is saved to
    `filename` instead (subject to the `override` flag).

    Parameters
    ----------
    pptx : Presentation or None
        Target presentation. When ``None``, the image is saved to disk.
    filename : str or Path
        Destination path used when `pptx` is ``None``.
    image : PIL.Image.Image
        The image to export.
    dpi : int, optional
        Resolution in dots per inch. Defaults to ``150``.
    fmt : str, optional
        PIL format string (e.g. ``"JPEG"``, ``"PNG"``). Defaults to
        :data:`DEFAULT_PIL_FMT`.
    override : bool, optional
        When ``True``, overwrite `filename` even if it already exists.
        Ignored when `pptx` is not ``None``. Defaults to ``False``.
    close : bool, optional
        Whether to close the image after export. Defaults to ``False``.
    title : str, optional
        Slide title inserted above the picture. Defaults to an empty string.
    **kwargs
        Additional keyword arguments forwarded to ``image.save``. A
        ``quality`` key defaults to ``95`` if not provided.
    """
    quality = kwargs.pop("quality", 95)
    if pptx is not None:
        if image.mode == "RGBA" and fmt.upper() in ("JPEG", "JPG"):
            image = image.convert("RGB")
        with io.BytesIO() as image_stream:
            image.save(image_stream, fmt, quality=quality, dpi=(dpi, dpi), **kwargs)
            slide, left, top = _insert_slide(pptx, title=title)
            slide.shapes.add_picture(image_stream, left, top)
    else:
        if override or not Path(filename).exists():
            image.save(filename, dpi=(dpi, dpi), **kwargs)
    if close:
        image.close()


def _insert_slide(pptx: Presentation, title: str = "") -> tuple[ty.Any, int, int]:
    """Append a new blank (or titled) slide and return positioning anchors.

    Selects the ``TITLE_AND_BLANK`` layout when a title is provided, or the
    ``BLANK`` layout otherwise. When a title is provided its font size is
    reduced to 20 pt, the title box is positioned in the top-left corner at
    full slide width, and the top anchor is set below the title box so that
    picture content does not overlap it.

    Parameters
    ----------
    pptx : Presentation
        The presentation to which the slide is appended.
    title : str, optional
        Title text placed at the top of the slide. Defaults to an empty
        string (no title).

    Returns
    -------
    slide : pptx.slide.Slide
        The newly created slide object.
    left : int
        Horizontal offset (in EMU) for subsequent shape placement.
    top : int
        Vertical offset (in EMU) for subsequent shape placement, accounting
        for the title box height when a title is present.
    """
    left = top = 0
    layout_index = SlideLayout.TITLE_AND_BLANK if title else SlideLayout.BLANK
    slide = pptx.slides.add_slide(pptx.slide_layouts[layout_index])
    if title:
        from pptx.util import Cm, Pt

        line_count = title.count("\n") + 1
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(20)
        slide.shapes.title.text = title
        height = slide.shapes.title.height // 2 * line_count
        slide.shapes.title.top = Cm(0)
        slide.shapes.title.left = Cm(0)
        slide.shapes.title.width = pptx.slide_width
        slide.shapes.title.height = height
        top = slide.shapes.title.top + height
    return slide, left, top
